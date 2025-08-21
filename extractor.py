# mcp_ai/extractor.py
from __future__ import annotations

import io
import os
import time
import logging
import shutil
from typing import Dict, Any, Optional, List, Tuple

# Unified, robust A/V pipeline lives here
from .av_utils import transcribe_av

log = logging.getLogger("extractor")

# ---------------------------
# Capability matrix (summary)
# ---------------------------
# PDF:       PyMuPDF -> pdfminer -> pdfplumber -> OCR (if allowed)
# DOCX:      docx2txt
# PPTX:      python-pptx
# XLSX:      openpyxl (values only)
# TEXT-ish:  txt, csv/tsv, log, md, json, xml, html (bs4), ini/cfg, yml/yaml, mod (if policy allows)
# IMAGES:    png, jpg/jpeg, tiff -> OCR cascade (easyocr -> tesseract -> rapidocr/paddle/azure)
# AUDIO:     wav, mp3, m4a, aac, flac, ogg, opus -> transcribe_av
# VIDEO:     mp4, mov, mkv, webm -> transcribe_av
# Archives:  metadata_only here (scanner may expand upstream)

# Common OS/system junk (also handled by scanner)
_SYSTEM_BASENAMES = {".DS_Store", "Thumbs.db", "desktop.ini"}

_AUDIO_EXTS = {".wav", ".mp3", ".m4a", ".aac", ".flac", ".ogg", ".opus"}
_VIDEO_EXTS = {".mp4", ".mov", ".mkv", ".webm"}
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp"}

# ---------------------------
# Helpers
# ---------------------------

def _ok(status: str, meta: Dict[str, Any], text: str = "", skipped_reason: Optional[str] = None) -> Dict[str, Any]:
    """
    Standard extractor result envelope:
      extraction_status: "ok" | "metadata_only" | "skipped"
      skipped_reason: str | None
      text: extracted text/transcript/plain text
      meta: engine details, pages, dpi, language, etc.
    """
    return {
        "extraction_status": status,
        "skipped_reason": skipped_reason,
        "text": text,
        "meta": meta,
    }

def _quick_skip(path: str, reason: str) -> Dict[str, Any]:
    return _ok("skipped", {"ext": os.path.splitext(path)[1].lower()}, "", reason)

def _size_bytes(path: str) -> int:
    try:
        return os.path.getsize(path)
    except Exception:
        return -1

def _read_text_small(path: str, max_mb: int = 25) -> Optional[str]:
    """Read small plaintext-like files with UTF-8→Latin-1 fallback. Returns None if too large / not text."""
    try:
        if _size_bytes(path) > max_mb * 1024 * 1024:
            return None
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        try:
            with open(path, "r", encoding="latin-1") as f:
                return f.read()
        except Exception:
            return None
    except Exception:
        return None

def _has_tool(cmd: str) -> bool:
    return shutil.which(cmd) is not None

def _maybe_detect_language(text: str, profile: Dict[str, Any]) -> Optional[str]:
    """
    Best-effort language detection for texty outputs (NOT used for A/V; av_utils already reports language).
    Enabled by profile['detect_language'] (False by default). Uses whatever is installed: langid or langdetect.
    """
    try:
        if not profile.get("detect_language"):
            return None
        t = (text or "").strip()
        if len(t) < int(profile.get("detect_language_min_chars", 50)):
            return None

        # Try langid (fast, no network)
        try:
            import langid  # type: ignore
            code, _ = langid.classify(t)
            return code or None
        except Exception:
            pass

        # Fallback: langdetect
        try:
            from langdetect import detect  # type: ignore
            code = detect(t)
            return code or None
        except Exception:
            pass
    except Exception as e:
        log.debug("language detection failed: %r", e)
    return None

# ---------------------------
# PDF extraction cascade
# ---------------------------

def _extract_pdf_text(path: str, profile: Dict[str, Any]) -> Dict[str, Any]:
    """
    Try PDF text extract via configured cascade:
      pymupdf -> pdfminer -> pdfplumber -> OCR (if profile.ocr true)
    If nothing yields text, return metadata_only (not ok).
    """
    cascade = list((profile.get("parse_pdf_order") or ["pymupdf", "pdfminer", "pdfplumber", "ocr"]))
    meta: Dict[str, Any] = {"ext": ".pdf", "engine": "none", "cascade": cascade}
    text = ""

    # 1) PyMuPDF
    if "pymupdf" in cascade:
        try:
            import fitz  # type: ignore
            t0 = time.time()
            doc = fitz.open(path)
            pages = doc.page_count
            parts = []
            for i in range(pages):
                parts.append(doc.load_page(i).get_text())
            text = "\n".join(parts)
            meta.update({"engine": "pymupdf", "pages": pages, "elapsed_s": round(time.time()-t0,3)})
            if text.strip():
                meta["language"] = _maybe_detect_language(text, profile)
                return _ok("ok", meta, text)
        except Exception as e:
            log.debug("pymupdf failed on %s: %r", path, e)

    # 2) pdfminer
    if "pdfminer" in cascade:
        try:
            t0 = time.time()
            from pdfminer.high_level import extract_text as pdfminer_extract_text  # type: ignore
            text = pdfminer_extract_text(path) or ""
            meta.update({"engine": "pdfminer", "elapsed_s": round(time.time()-t0,3)})
            if text.strip():
                meta["language"] = _maybe_detect_language(text, profile)
                return _ok("ok", meta, text)
        except Exception as e:
            log.debug("pdfminer failed on %s: %r", path, e)

    # 3) pdfplumber
    if "pdfplumber" in cascade:
        try:
            import pdfplumber  # type: ignore
            t0 = time.time()
            parts = []
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    parts.append(page.extract_text() or "")
            text = "\n".join([p for p in parts if p])
            meta.update({"engine": "pdfplumber", "elapsed_s": round(time.time()-t0,3)})
            if text.strip():
                meta["language"] = _maybe_detect_language(text, profile)
                return _ok("ok", meta, text)
        except Exception as e:
            log.debug("pdfplumber failed on %s: %r", path, e)

    # 4) OCR (if allowed)
    if "ocr" in cascade and profile.get("ocr"):
        return _extract_pdf_ocr(path, profile, meta_base=meta, max_pages=int(profile.get("ocr_max_pages", 10)))

    # No crash; just no text anywhere + OCR not allowed
    return _ok("metadata_only", meta, "", "no extractable text (no text layer / OCR disabled)")

def _extract_pdf_ocr(path: str, profile: Dict[str, Any], meta_base: Dict[str, Any], max_pages: int) -> Dict[str, Any]:
    """
    Render each page to an image (via PyMuPDF) and OCR by the image OCR cascade.
    Timeboxed by profile['timebox_ocr_s'].
    """
    try:
        import fitz  # type: ignore
    except Exception:
        return _ok("metadata_only", dict(meta_base, engine="ocr"), "", "pymupdf not available for OCR")

    lang = profile.get("ocr_lang", "eng")
    dpi = int(profile.get("ocr_dpi", 144))
    engines = list(profile.get("ocr_engine_order") or ["easyocr", "tesseract", "rapidocr", "paddle", "azure"])
    timebox_s = int(profile.get("timebox_ocr_s", 45))

    t0_all = time.time()
    text_parts: List[str] = []
    pages_ocrd = 0

    try:
        doc = fitz.open(path)
        pages = min(doc.page_count, max_pages)
    except Exception as e:
        log.debug("pymupdf open for OCR failed on %s: %r", path, e)
        return _ok("metadata_only", dict(meta_base, engine="ocr"), "", "pymupdf open failed")

    for i in range(pages):
        if time.time() - t0_all > timebox_s:
            break
        try:
            page = doc.load_page(i)
            pix = page.get_pixmap(dpi=dpi)
            img_bytes = pix.tobytes("png")
            t_text, engine_used = _ocr_bytes(img_bytes, engines, lang, timebox_s=max(5, timebox_s//2))
            if t_text:
                text_parts.append(t_text)
            pages_ocrd += 1
        except Exception as e:
            log.debug("OCR page %s failed: %r", i, e)

    meta = dict(meta_base)
    meta.update({
        "engine": "ocr",
        "ocr_pages": pages_ocrd,
        "dpi": dpi,
        "lang": lang,
        "elapsed_s": round(time.time() - t0_all, 3),
        "ocr_engine_order": engines,
    })
    text = "\n".join(text_parts).strip()
    if text:
        meta["language"] = _maybe_detect_language(text, profile)
        return _ok("ok", meta, text)
    return _ok("metadata_only", meta, "", "ocr produced no text")

# ---------------------------
# Image OCR cascade
# ---------------------------

def _ocr_bytes(img_bytes: bytes, engines: List[str], lang: str, timebox_s: int = 30) -> Tuple[str, str]:
    """
    Try engines in order; return (text, engine_used_or_reason).
    """
    t0 = time.time()

    # EASYOCR
    if "easyocr" in engines:
        try:
            import easyocr  # type: ignore
            reader = easyocr.Reader([lang], gpu=False)
            res = reader.readtext(img_bytes, detail=0)
            text = "\n".join([r for r in res if r])
            if text.strip():
                return text, "easyocr"
        except Exception as e:
            log.debug("easyocr failed: %r", e)
        if time.time() - t0 > timebox_s:
            return "", "easyocr-timeout"

    # TESSERACT
    if "tesseract" in engines:
        try:
            from PIL import Image  # type: ignore
            import pytesseract  # type: ignore
            img = Image.open(io.BytesIO(img_bytes))
            text = pytesseract.image_to_string(img, lang=lang)
            if text.strip():
                return text, "tesseract"
        except Exception as e:
            log.debug("tesseract failed: %r", e)
        if time.time() - t0 > timebox_s:
            return "", "tesseract-timeout"

    # RAPIDOCR (optional)
    if "rapidocr" in engines:
        try:
            from rapidocr_onnxruntime import RapidOCR  # type: ignore
            ocr = RapidOCR()
            res, _ = ocr(io.BytesIO(img_bytes))
            text = "\n".join([x[1] for x in (res or []) if x and len(x) > 1])
            if text.strip():
                return text, "rapidocr"
        except Exception as e:
            log.debug("rapidocr failed: %r", e)

    # PADDLE (optional placeholder)
    if "paddle" in engines:
        try:
            # If you wire PaddleOCR, add it here.
            pass
        except Exception as e:
            log.debug("paddle failed: %r", e)

    # AZURE (disabled for local-only; the router would call cloud if allowed)
    if "azure" in engines:
        return "", "azure-disabled"

    return "", "none"

# ---------------------------
# Office/markup handlers
# ---------------------------

def _extract_docx(path: str) -> Dict[str, Any]:
    try:
        import docx2txt  # type: ignore
        t0 = time.time()
        txt = docx2txt.process(path) or ""
        meta = {"ext": ".docx", "engine": "docx2txt", "elapsed_s": round(time.time()-t0,3)}
        if txt.strip():
            meta["language"] = _maybe_detect_language(txt, {})
        return _ok("ok", meta, txt)
    except Exception as e:
        log.debug("docx2txt failed on %s: %r", path, e)
        return _ok("metadata_only", {"ext": ".docx", "engine": "none"}, "", "docx2txt not available or failed")

def _extract_pptx(path: str) -> Dict[str, Any]:
    try:
        from pptx import Presentation  # type: ignore
        t0 = time.time()
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        txt = "\n".join(parts)
        meta = {"ext": ".pptx", "engine": "python-pptx", "elapsed_s": round(time.time()-t0,3)}
        if txt.strip():
            meta["language"] = _maybe_detect_language(txt, {})
        return _ok("ok", meta, txt)
    except Exception as e:
        log.debug("python-pptx failed on %s: %r", path, e)
        return _ok("metadata_only", {"ext": ".pptx", "engine": "none"}, "", "python-pptx not available or failed")

def _extract_xlsx(path: str) -> Dict[str, Any]:
    try:
        from openpyxl import load_workbook  # type: ignore
        t0 = time.time()
        wb = load_workbook(path, read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                vals = [str(v) for v in row if v is not None]
                if vals:
                    parts.append(", ".join(vals))
        txt = "\n".join(parts)
        meta = {"ext": ".xlsx", "engine": "openpyxl", "elapsed_s": round(time.time()-t0,3)}
        if txt.strip():
            meta["language"] = _maybe_detect_language(txt, {})
        return _ok("ok", meta, txt)
    except Exception as e:
        log.debug("openpyxl failed on %s: %r", path, e)
        return _ok("metadata_only", {"ext": ".xlsx", "engine": "none"}, "", "openpyxl not available or failed")

def _extract_html(path: str, profile: Dict[str, Any]) -> Dict[str, Any]:
    # Prefer BeautifulSoup if available
    try:
        from bs4 import BeautifulSoup  # type: ignore
        t0 = time.time()
        html = _read_text_small(path, max_mb=10) or ""
        soup = BeautifulSoup(html, "html.parser")
        txt = (soup.get_text("\n") or "").strip()
        meta = {"ext": ".html", "engine": "bs4", "elapsed_s": round(time.time()-t0,3)}
        if txt:
            meta["language"] = _maybe_detect_language(txt, profile)
            return _ok("ok", meta, txt)
        return _ok("metadata_only", meta, "", "no extractable text")
    except Exception:
        txt = _read_text_small(path, max_mb=10)
        if txt is not None:
            txt = (txt or "").strip()
            meta = {"ext": ".html", "engine": "text"}
            if txt:
                meta["language"] = _maybe_detect_language(txt, profile)
                return _ok("ok", meta, txt)
            return _ok("metadata_only", meta, "", "no extractable text")
        return _ok("metadata_only", {"ext": ".html", "engine": "none"}, "", "html too large or bs4 not available")

# ---------------------------
# Image handler (+EXIF)
# ---------------------------

def _extract_image(path: str, profile: Dict[str, Any]) -> Dict[str, Any]:
    lang = profile.get("ocr_lang", "eng")
    engines = list(profile.get("ocr_engine_order") or ["easyocr", "tesseract", "rapidocr", "paddle", "azure"])
    timebox_s = int(profile.get("timebox_ocr_s", 45))
    collect_exif = bool(profile.get("collect_exif", False))

    meta: Dict[str, Any] = {"ext": os.path.splitext(path)[1].lower(), "engine": "image", "ocr_engine_order": engines}
    exif: Dict[str, Any] = {}

    img_bytes = None
    try:
        with open(path, "rb") as f:
            img_bytes = f.read()
    except Exception as e:
        log.debug("image read failed: %r", e)
        return _ok("metadata_only", meta, "", "image read failed")

    # OCR (if profile.ocr true)
    text = ""
    if profile.get("ocr"):
        t_text, used = _ocr_bytes(img_bytes, engines, lang, timebox_s=timebox_s)
        text = (t_text or "").strip()
        meta["engine"] = used

    # EXIF
    if collect_exif:
        try:
            from PIL import Image, ExifTags  # type: ignore
            im = Image.open(io.BytesIO(img_bytes))
            raw = getattr(im, "_getexif", lambda: None)() or {}
            inv = {v: k for k, v in ExifTags.TAGS.items()}
            wanted = ["DateTimeOriginal", "Model", "Make", "Orientation", "XResolution", "YResolution"]
            for k, v in (raw.items() if isinstance(raw, dict) else []):
                tag = ExifTags.TAGS.get(k) or inv.get(k) or str(k)
                if tag in wanted:
                    exif[tag] = v
            if exif:
                meta["exif"] = exif
        except Exception as e:
            log.debug("exif read failed: %r", e)

    if text:
        meta["language"] = _maybe_detect_language(text, profile)
        return _ok("ok", meta, text)

    # If OCR enabled but produced no text, mark metadata_only (this fixes the PNG issue you saw)
    if profile.get("ocr"):
        return _ok("metadata_only", meta, "", "no text recognized")

    # If OCR disabled, still metadata_only since we didn't read text
    return _ok("metadata_only", meta, "", "ocr disabled for images")

# ---------------------------
# Main dispatch
# ---------------------------

def extract_content(
    path: str,
    profile: Dict[str, Any],
    type_policies: Dict[str, Any],
    cfg: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    """
    Unified dispatcher. cfg is optional (kept for backward compatibility).
    If provided, A/V will use ASR settings from cfg.ai.local.faster_whisper.
    """
    base = os.path.basename(path)
    if base in _SYSTEM_BASENAMES:
        return _quick_skip(path, "system file")

    ext = os.path.splitext(path)[1].lower()

    # 0) Hard skip types (policy)
    for s in (type_policies.get("skip") or []):
        if ext == s.lower():
            return _quick_skip(path, f"skip policy for {ext}")

    # 0.5) Metadata-only (policy)
    for s in (type_policies.get("metadata_only") or []):
        if ext == s.lower():
            return _ok("metadata_only", {"ext": ext, "engine": "none"}, "", f"policy: metadata_only for {ext}")

    # 1) PDFs
    if ext == ".pdf":
        return _extract_pdf_text(path, profile)

    # 2) Images (with optional OCR + EXIF)
    if ext in _IMAGE_EXTS:
        return _extract_image(path, profile)

    # 3) Audio/Video → single, robust path (transcribe_av)
    if ext in _AUDIO_EXTS or ext in _VIDEO_EXTS:
        try:
            # cfg may be None; transcribe_av handles defaults internally
            return transcribe_av(path, profile, cfg or {})
        except Exception as e:
            log.warning("transcribe_av failed on %s: %r", path, e)
            return _ok("metadata_only", {"ext": ext, "engine": "none"}, "", "transcribe_av failed")

    # 4) Office + structured
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".xlsx":
        return _extract_xlsx(path)
    if ext in (".html", ".htm"):
        return _extract_html(path, profile)

    # 5) Plaintext-ish (incl. .mod if NOT policy-blocked)
    if ext in (".txt", ".csv", ".tsv", ".log", ".md", ".json", ".xml", ".ini", ".cfg", ".yml", ".yaml", ".mod"):
        # If policy says metadata_only for this ext, honor it
        if ext in [e.lower() for e in (type_policies.get("metadata_only") or [])]:
            return _ok("metadata_only", {"ext": ext, "engine": "none"}, "", f"policy: metadata_only for {ext}")

        txt = _read_text_small(path, max_mb=25)
        if txt is not None:
            txt_stripped = (txt or "").strip()
            meta = {"ext": ext, "engine": "text"}
            if txt_stripped:
                meta["language"] = _maybe_detect_language(txt_stripped, profile)
                return _ok("ok", meta, txt)
            return _ok("metadata_only", meta, "", "no extractable text")
        # try_text_then_metadata_only fallback (e.g., very large .json/.log)
        if ext in [e.lower() for e in (type_policies.get("try_text_then_metadata_only") or [])]:
            return _ok("metadata_only", {"ext": ext, "engine": "none"}, "", "policy: try_text_then_metadata_only → metadata_only")
        return _ok("metadata_only", {"ext": ext, "engine": "none"}, "", "too large for plaintext ingest")

    # 6) Archives - prefer upstream expansion. Provide metadata_only here.
    if ext in (".zip", ".tar", ".gz", ".tgz", ".rar", ".7z"):
        return _ok("metadata_only", {"ext": ext, "engine": "none"}, "", "archive container (expand upstream)")

    # 7) Metadata-only fallback (unknown type)
    meta = {"ext": ext, "engine": "none"}
    return _ok("metadata_only", meta, "", "unknown type")